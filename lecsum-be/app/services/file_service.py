# PDF / PPT 처리 로직
import os
import shutil
import uuid

from fastapi import UploadFile, HTTPException
from sqlalchemy.orm import Session

# LangChain
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain_core.documents import Document


from app.core.llm_client import summary_chain, keyword_chain, top_sentence_chain

from pptx import Presentation

from app.crud import file_crud
from app.core.enums import ChromaDB, CHROMA_PERSIST_DIR
from app.db.file_schemas import (DocumentSummaryItem, DocumentSummaryDetail)

def extract_text_from_pptx(file_path: str) -> str:
    """PPTX 파일에서 텍스트 추출"""
    prs = Presentation(file_path)
    texts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

    return "\n".join(texts)

async def register_document(db: Session, file: UploadFile, summary_type: str):
    """
    PDF / PPT 업로드 → 요약 생성 → VectorDB 저장 → MySQL 저장
    Returns:
        dict: {"id": 문서 ID, "summary": 요약}
    """
    file_uuid = str(uuid.uuid4())
    filename_lower = file.filename.lower()
    temp_file_path = f"temp_{file.filename}"

    try:
        # 임시 파일 저장
        with open(temp_file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # 파일 유형별 텍스트 로드
        if filename_lower.endswith(".pdf"):
            loader = PyPDFLoader(temp_file_path)
            docs = loader.load()

        elif filename_lower.endswith(".pptx"):
            text = extract_text_from_pptx(temp_file_path)
            docs = [Document(page_content=text, metadata={})]

        else:
            raise HTTPException(
                status_code=400,
                detail="지원하지 않는 파일 형식입니다. (PDF / PPTX만 가능)"
            )

        full_text = "\n".join(doc.page_content for doc in docs)

        # 요약 생성
        summary_result = summary_chain.invoke({
            "context": full_text,
            "summary_type": summary_type
        })

        summary = (
            summary_result.content
            if hasattr(summary_result, "content")
            else str(summary_result)
        )

        keyword_result = keyword_chain.invoke({"context": full_text})

        # 키워드 결과 정규화 (LLM 출력 형태 다양성 대응)
        if hasattr(keyword_result, "content"):
            raw_keywords = keyword_result.content
        else:
            raw_keywords = str(keyword_result)

        # 문자열 → 리스트 변환 (쉼표/줄바꿈 기준)
        keywords = [
            k.strip()
            for k in raw_keywords.replace("\n", ",").split(",")
            if k.strip()
        ]

        top_sentences = top_sentence_chain.invoke({
            "context": full_text,
            "k": 5
        })

        keywords_str = ", ".join(keywords)

        # 문서 청킹
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        splits = splitter.split_documents(docs)

        for doc in splits:
            doc.metadata.update({
                "document_id": doc.id,
                "document_uuid": file_uuid,
                "filename": file.filename,
                "keywords": keywords_str,
            })

        # ChromaDB 저장
        vectorstore = Chroma(
            collection_name=ChromaDB.COLLECTION_NAME.value,
            embedding_function=OpenAIEmbeddings(
                model="text-embedding-3-small"
            ),
            persist_directory=CHROMA_PERSIST_DIR,
        )
        vectorstore.add_documents(splits)
        
                # 추가 메타데이터 계산
        concept_count = summary.count("###")
        keyword_count = len(keywords)
        word_count = len(summary.split())
        review_time_min = max(1, word_count // 50)

        # MySQL 저장
        document = file_crud.create_document(
            db=db,
            uuid=file_uuid,
            name=file.filename,
            summary=summary,
            keywords=keywords_str,
            concept_cnt=concept_count,
            keyword_cnt=keyword_count,
            review_time=review_time_min
        )

        return DocumentSummaryDetail(
            id=document.id,
            uuid=document.uuid,
            name=document.name,
            summary=document.summary,
            keywords=document.keywords.split(", ") if document.keywords else [],
            concept_cnt=concept_count,
            keyword_cnt=keyword_count,
            review_time=review_time_min,
            created_at=document.created_at,
        )

    finally:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)

def list_documents(db: Session, limit: int, offset: int):
    """
    업로드된 문서 요약 목록 조회
    """
    document_files = file_crud.list_documents(
        db=db,
        limit=limit,
        offset=offset,
    )

    results: list[DocumentSummaryItem] = []

    for document in document_files:
        results.append(
            DocumentSummaryItem(
                id=document.id,
                name=document.name,
                summary=document.summary,
                keywords=document.keywords.split(", ") if document.keywords else [],
                concept_cnt=document.concept_cnt,
                keyword_cnt=document.keyword_cnt,
                review_time=document.review_time,
                created_at=document.created_at,
            )
        )

    return results

def get_document_detail(db: Session, document_id: int):
    """
    특정 문서 요약 상세 조회
    """
    document = file_crud.get_document_by_id(db, document_id)

    if not document:
        raise HTTPException(
            status_code=404,
            detail="문서를 찾을 수 없습니다."
        )

    return DocumentSummaryDetail(
        id=document.id,
        uuid=document.uuid,
        name=document.name,
        summary=document.summary,
        keywords=document.keywords.split(", ") if document.keywords else [],
        created_at=document.created_at,
    )
